import base64
import csv
import io
import json
import os
import re

import anthropic
import fitz  # PyMuPDF
import openpyxl
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from openpyxl.styles import Alignment, Font, PatternFill
from PIL import Image
from pydantic import BaseModel

load_dotenv()

app = FastAPI(title="Chart Data Extractor API")

ALLOWED_ORIGINS = [
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    # Add your Vercel URL here after deploying, e.g.:
    # "https://chart-data-extractor.vercel.app",
    *[o.strip() for o in os.getenv("ALLOWED_ORIGINS", "").split(",") if o.strip()],
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

claude = anthropic.AsyncAnthropic()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def compress_image(image_bytes: bytes, max_side: int = 1400) -> tuple[str, str]:
    """Resize & convert to JPEG; return (base64_string, media_type)."""
    img = Image.open(io.BytesIO(image_bytes))
    if img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGB")
    w, h = img.size
    if max(w, h) > max_side:
        ratio = max_side / max(w, h)
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85, optimize=True)
    return base64.b64encode(buf.getvalue()).decode(), "image/jpeg"


def safe_sheet_name(name: str, existing: list[str]) -> str:
    """Return a unique, Excel-safe sheet name (max 31 chars)."""
    base = re.sub(r"[\\/*?:\[\]]", "_", name)[:28] or "Chart"
    candidate = base
    suffix = 1
    while candidate in existing:
        candidate = f"{base}_{suffix}"
        suffix += 1
    return candidate


EXTRACTION_PROMPT = """Analyze this image and extract every chart or graph you can find.

For EACH chart:
1. Identify chart type: bar, horizontal_bar, stacked_bar, stacked_horizontal_bar,
   grouped_bar, pie, donut, line, area, scatter
2. Extract the title (if visible)
3. Extract ALL data points with exact labels and numeric values
4. Identify legend / series names for multi-series charts
5. Note the unit (%, $, count, etc.) if shown

Return ONLY a valid JSON object — no markdown, no explanation:

{
  "has_charts": true,
  "confidence": 0.95,
  "charts": [
    {
      "id": 1,
      "type": "stacked_horizontal_bar",
      "title": "Survey Results 2024",
      "unit": "%",
      "series": ["Strongly Agree", "Agree", "Neutral", "Disagree", "Strongly Disagree"],
      "data": [
        {
          "label": "I feel valued at work",
          "values": {
            "Strongly Agree": 32,
            "Agree": 41,
            "Neutral": 15,
            "Disagree": 8,
            "Strongly Disagree": 4
          }
        }
      ]
    }
  ]
}

For single-series charts use "value" (number), not "values":
  "data": [{"label": "January", "value": 42.5}]

For pie / donut charts use "value" per slice.

If no charts are found:
{"has_charts": false, "confidence": 1.0, "charts": []}

Return ONLY the JSON object."""


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.get("/health")
async def health():
    return {"status": "ok"}


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    """Accept PDF, PNG, JPG, PPTX. Return list of page images (base64 JPEG)."""
    content = await file.read()
    filename = (file.filename or "").lower()
    images: list[dict] = []

    # ---- PDF ----
    if filename.endswith(".pdf"):
        doc = fitz.open(stream=content, filetype="pdf")
        for i in range(min(len(doc), 30)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            b64, mt = compress_image(pix.tobytes("png"))
            images.append({"page": i + 1, "data": b64, "media_type": mt})
        doc.close()

    # ---- Raster images ----
    elif filename.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")):
        b64, mt = compress_image(content)
        images.append({"page": 1, "data": b64, "media_type": mt})

    # ---- PPTX: extract embedded picture shapes ----
    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = Presentation(io.BytesIO(content))
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_images = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            b64, mt = compress_image(shape.image.blob)
                            slide_images.append(
                                {"page": slide_num, "data": b64, "media_type": mt}
                            )
                        except Exception:
                            pass
                # If no picture shapes, render a blank placeholder so user knows
                if not slide_images:
                    slide_images.append(
                        {"page": slide_num, "data": "", "media_type": "image/jpeg",
                         "placeholder": True}
                    )
                images.extend(slide_images[:1])  # one image per slide
        except ImportError:
            raise HTTPException(
                400,
                "python-pptx is required for PPTX support. "
                "Run: pip install python-pptx",
            )
        except Exception as e:
            raise HTTPException(400, f"PPTX processing error: {e}")

    else:
        raise HTTPException(
            400,
            f"Unsupported file type '{filename}'. "
            "Supported: PDF, PNG, JPG, JPEG, WEBP, PPTX",
        )

    if not images:
        raise HTTPException(400, "Could not extract any images from the file.")

    return {
        "filename": file.filename,
        "total_pages": len(images),
        "images": images,
    }


class ExtractRequest(BaseModel):
    image: str
    media_type: str = "image/jpeg"


@app.post("/api/extract")
async def extract_chart(request: ExtractRequest):
    """Send a page image to Claude and return structured chart data."""
    if not request.image:
        raise HTTPException(400, "Empty image data")

    try:
        async with claude.messages.stream(
            model="claude-opus-4-6",
            max_tokens=8192,
            thinking={"type": "adaptive"},
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": request.media_type,
                                "data": request.image,
                            },
                        },
                        {"type": "text", "text": EXTRACTION_PROMPT},
                    ],
                }
            ],
        ) as stream:
            final = await stream.get_final_message()

        text = next(
            (b.text for b in final.content if b.type == "text"), "{}"
        ).strip()

        # Strip markdown fences if present
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)

        return json.loads(text)

    except json.JSONDecodeError:
        # Try to salvage JSON from the response
        match = re.search(r"\{.*\}", text, re.DOTALL)
        if match:
            try:
                return json.loads(match.group())
            except Exception:
                pass
        return {"has_charts": False, "confidence": 0.0, "charts": []}
    except Exception as e:
        raise HTTPException(500, f"Extraction failed: {e}")


class ExportRequest(BaseModel):
    charts: list[dict]
    format: str = "xlsx"
    filename: str = "chart_data"


@app.post("/api/export")
async def export_data(request: ExportRequest):
    """Convert extracted chart data to xlsx / csv / json download."""
    charts = request.charts
    if not charts:
        raise HTTPException(400, "No charts provided for export.")

    fmt = request.format.lower()
    fn = re.sub(r"[^\w\-_]", "_", request.filename) or "chart_data"

    # ---- JSON ----
    if fmt == "json":
        payload = json.dumps(charts, indent=2, ensure_ascii=False).encode()
        return StreamingResponse(
            io.BytesIO(payload),
            media_type="application/json",
            headers={"Content-Disposition": f'attachment; filename="{fn}.json"'},
        )

    # ---- CSV ----
    if fmt == "csv":
        buf = io.StringIO()
        w = csv.writer(buf)
        for i, chart in enumerate(charts):
            if i:
                w.writerow([])
            title = chart.get("title") or f"Chart {i + 1}"
            ctype = chart.get("type", "unknown")
            unit = chart.get("unit", "")
            series = chart.get("series") or []
            data = chart.get("data") or []

            w.writerow([f"Chart: {title}", f"Type: {ctype}", f"Unit: {unit}"])
            w.writerow([])

            if series:
                w.writerow(["Category"] + series)
                for row in data:
                    vals = row.get("values") or {}
                    w.writerow([row.get("label", "")] + [vals.get(s, "") for s in series])
            else:
                label = f"Value ({unit})" if unit else "Value"
                w.writerow(["Category", label])
                for row in data:
                    w.writerow([row.get("label", ""), row.get("value", "")])

        csv_bytes = buf.getvalue().encode("utf-8-sig")  # BOM for Excel compat
        return StreamingResponse(
            io.BytesIO(csv_bytes),
            media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="{fn}.csv"'},
        )

    # ---- XLSX ----
    if fmt == "xlsx":
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default empty sheet

        hdr_font = Font(bold=True, color="FFFFFF", size=10)
        hdr_fill = PatternFill("solid", fgColor="1E3A5F")
        meta_font = Font(bold=True, color="374151", size=9)

        existing_sheets: list[str] = []

        for i, chart in enumerate(charts):
            title = chart.get("title") or f"Chart {i + 1}"
            ctype = chart.get("type", "unknown")
            unit = chart.get("unit", "")
            series = chart.get("series") or []
            data = chart.get("data") or []

            sname = safe_sheet_name(title, existing_sheets)
            existing_sheets.append(sname)
            ws = wb.create_sheet(title=sname)

            # Metadata block
            ws["A1"], ws["B1"] = "Chart", title
            ws["A2"], ws["B2"] = "Type", ctype.replace("_", " ").title()
            ws["A3"], ws["B3"] = "Unit", unit
            for r in range(1, 4):
                ws.cell(r, 1).font = meta_font
            ws.append([])

            # Header row
            if series:
                headers = ["Category"] + series
            else:
                lbl = f"Value ({unit})" if unit else "Value"
                headers = ["Category", lbl]

            ws.append(headers)
            hdr_row = ws.max_row
            for col in range(1, len(headers) + 1):
                cell = ws.cell(hdr_row, col)
                cell.font = hdr_font
                cell.fill = hdr_fill
                cell.alignment = Alignment(horizontal="center")

            # Data rows
            for row in data:
                if series:
                    vals = row.get("values") or {}
                    ws.append([row.get("label", "")] + [vals.get(s, "") for s in series])
                else:
                    ws.append([row.get("label", ""), row.get("value", "")])

            # Auto-width columns
            for col_cells in ws.columns:
                width = max(
                    (len(str(c.value or "")) for c in col_cells), default=8
                )
                ws.column_dimensions[col_cells[0].column_letter].width = min(width + 3, 60)

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)
        return StreamingResponse(
            out,
            media_type=(
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ),
            headers={"Content-Disposition": f'attachment; filename="{fn}.xlsx"'},
        )

    raise HTTPException(400, f"Unknown format '{fmt}'. Use xlsx, csv, or json.")
